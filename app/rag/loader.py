"""
Document loaders — turn one uploaded/notes file into LangChain Documents.

We support PDF, DOCX, Markdown and plain text. For PDF we try pdfplumber first
(generally best text quality on lecture notes) and fall back to pypdf.

Each chunk's metadata records where it came from so we can show sources to the
user and (later) delete/re-ingest a single file.
"""
from __future__ import annotations

from pathlib import Path
from typing import List

from langchain_core.documents import Document

from app.core.logging import get_logger
from app.rag.text_utils import clean_note_text, sanitize_metadata

log = get_logger(__name__)

SUPPORTED_EXTS = {".pdf", ".docx", ".md", ".txt", ".markdown", ".pptx"}


def _load_pdf(path: Path) -> List[Document]:
    """Load a PDF using pdfplumber, fall back to pypdf."""
    # pdfplumber
    try:
        from langchain_community.document_loaders import PDFPlumberLoader

        docs = PDFPlumberLoader(str(path)).load()
        if docs and any(d.page_content.strip() for d in docs):
            return docs
    except Exception as e:  # noqa: BLE001
        log.warning(f"pdfplumber failed on {path.name}: {e}")

    # fallback: pypdf
    from langchain_community.document_loaders import PyPDFLoader

    return PyPDFLoader(str(path)).load()


def _load_docx(path: Path) -> List[Document]:
    from langchain_community.document_loaders import Docx2txtLoader

    try:
        return Docx2txtLoader(str(path)).load()
    except Exception:  # noqa: BLE001
        # python-docx fallback
        import docx  # type: ignore

        d = docx.Document(str(path))
        text = "\n".join(p.text for p in d.paragraphs)
        return [Document(page_content=text, metadata={"source": str(path)})]


def _load_text(path: Path) -> List[Document]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    return [Document(page_content=text, metadata={"source": str(path)})]


def _load_pptx(path: Path) -> List[Document]:
    """Load a PowerPoint (.pptx) file — one Document per slide."""
    from pptx import Presentation  # python-pptx

    prs = Presentation(str(path))
    docs: List[Document] = []
    for i, slide in enumerate(prs.slides, start=1):
        # Collect every text frame on the slide: title, body, tables, notes.
        texts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        texts.append(row_text)
        # Speaker notes are gold for studying — include them.
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"(Speaker notes: {notes})")

        content = "\n".join(texts).strip()
        if content:
            docs.append(
                Document(
                    page_content=content,
                    metadata={"source": str(path), "page": i, "slide": i},
                )
            )
    if not docs:
        docs = [Document(page_content="", metadata={"source": str(path)})]
    return docs


def _load_markdown(path: Path) -> List[Document]:
    """Load markdown, strip formatting to plain text for cleaner embeddings."""
    import markdown as md  # type: ignore
    import re

    raw = path.read_text(encoding="utf-8", errors="ignore")
    html = md.markdown(raw)
    # strip HTML tags to get readable plain text
    text = re.sub(r"<[^>]+>", "", html)
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    return [Document(page_content=text, metadata={"source": str(path)})]


def load_single_file(path: Path) -> List[Document]:
    """Load a single file based on its extension. Raises ValueError on unknown."""
    ext = path.suffix.lower()
    if ext == ".pdf":
        docs = _load_pdf(path)
    elif ext == ".docx":
        docs = _load_docx(path)
    elif ext in (".md", ".markdown"):
        docs = _load_markdown(path)
    elif ext == ".pptx":
        docs = _load_pptx(path)
    elif ext == ".txt":
        docs = _load_text(path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    # Tag every document with the source filename so we can show it in the UI.
    source_name = path.name
    for d in docs:
        d.page_content = clean_note_text(d.page_content)
        d.metadata.setdefault("source", str(path))
        d.metadata["source_name"] = source_name
        d.metadata["file_type"] = ext.lstrip(".")
        d.metadata = sanitize_metadata(d.metadata)
    log.info(f"Loaded {len(docs)} section(s) from {source_name}")
    return [d for d in docs if d.page_content.strip()]


def load_directory(directory: Path) -> List[Document]:
    """Recursively load every supported file in `directory`."""
    if not directory.exists():
        log.warning(f"Notes directory does not exist yet: {directory}")
        return []

    all_docs: List[Document] = []
    files = sorted(
        p for p in directory.rglob("*") if p.is_file() and p.suffix.lower() in SUPPORTED_EXTS
    )
    log.info(f"Found {len(files)} supported note file(s) in {directory}")
    for p in files:
        try:
            all_docs.extend(load_single_file(p))
        except Exception as e:  # noqa: BLE001
            log.error(f"Failed to load {p}: {e}")
    return all_docs
